import json
import os
from pptx import Presentation
from pptx.util import Inches
from glob import glob

# 경로 설정
TEMPLATE_PATH = "template.pptx"
HYMN_DIR = "hymns"
SCRIPTURE_DIR = "scriptures"
RESPONSIVE_DIR = "responsive_readings"
CONFIG_PATH = "configuration.json"
OUTPUT_PATH = "output.pptx"

def load_configuration(path):
    with open(path, 'r', encoding='utf-8') as f:
        return json.load(f)

def append_slides_from(pptx_path, target_pres):
    if not os.path.exists(pptx_path):
        print(f"[경고] {pptx_path} 파일 없음.")
        return
    src = Presentation(pptx_path)
    for slide in src.slides:
        target_slide = target_pres.slides.add_slide(target_pres.slide_layouts[6])  # 빈 레이아웃
        for shape in slide.shapes:
            el = shape.element
            new_el = el.clone()
            target_slide.shapes._spTree.insert_element_before(new_el, 'p:extLst')

def add_title_slide(pres, title, date):
    slide = pres.slides.add_slide(pres.slide_layouts[0])  # 제목 슬라이드
    slide.shapes.title.text = title
    slide.placeholders[1].text = date  # 부제목에 날짜

def main():
    config = load_configuration(CONFIG_PATH)
    pres = Presentation(TEMPLATE_PATH)

    # 제목 슬라이드
    add_title_slide(pres, config["제목"], config["날짜"])

    # 찬송가
    for key in ["찬송1", "찬송2", "찬송3"]:
        hymn_num = config.get(key)
        if hymn_num:
            path = os.path.join(HYMN_DIR, f"[새찬송가] {hymn_num.zfill(3)}장.pptx")
            append_slides_from(path, pres)

    # 교독문
    reading_num = config.get("교독문")
    if reading_num:
        path = os.path.join(RESPONSIVE_DIR, f"교독문 {reading_num}.pptx")
        append_slides_from(path, pres)

    # 성경봉독 (구약 / 신약)
    bible = config.get("성경봉독", {})
    for category, text in bible.items():
        # category = "구약" or "신약"
        dir_map = {"구약": "OT", "신약": "NT"}
        scripture_path = os.path.join(SCRIPTURE_DIR, dir_map[category], f"{text}.pptx")
        append_slides_from(scripture_path, pres)

    # 기도자 슬라이드
    prayer = config.get("기도")
    if prayer:
        slide = pres.slides.add_slide(pres.slide_layouts[1])
        slide.shapes.title.text = "대표 기도"
        slide.placeholders[1].text = prayer

    # 저장
    pres.save(OUTPUT_PATH)
    print(f"[완료] {OUTPUT_PATH} 저장됨.")

if __name__ == "__main__":
    main()
