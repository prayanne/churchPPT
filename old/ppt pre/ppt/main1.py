"""
demo_ppt_maker.py
-----------------
• 기존 템플릿(.pptx) 유무에 따라 새 프레젠테이션을 만들고
• 예시 슬라이드 2장을 채워 저장합니다.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
import sys
from pathlib import Path

def create_ppt(template: str | None = None,
               output: str = "demo_output.pptx") -> None:
    """
    template : 기존 .pptx 양식 경로 (없으면 기본 테마)
    output   : 저장할 파일명
    """
    prs = Presentation(template) if template else Presentation()

    # --- 1. 제목 슬라이드 ---
    title_slide_layout = prs.slide_layouts[0]          # 첫 레이아웃(제목+부제목)
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = "데모 프레젠테이션"
    slide.placeholders[1].text = "python-pptx 기본 템플릿 예시"

    # --- 2. 내용(글머리) 슬라이드 ---
    bullet_layout = prs.slide_layouts[1]               # 제목+본문(글머리) 레이아웃
    slide2 = prs.slides.add_slide(bullet_layout)
    slide2.shapes.title.text = "📌 오늘의 Agenda"

    body = slide2.shapes.placeholders[1].text_frame
    for txt in ["프로젝트 개요", "요구 사항", "다음 단계"]:
        p = body.add_paragraph() if body.text else body.paragraphs[0]
        p.text = txt
        p.level = 0                                    # 글머리 수준(0 = 가장 바깥)
        p.font.size = Pt(18)

    prs.save(output)
    print(f"[완료] '{output}' 생성!")

if __name__ == "__main__":
    tpl = sys.argv[1] if len(sys.argv) > 2 else None   # python demo_ppt_maker.py 템플릿.pptx 출력명.pptx
    out = sys.argv[2] if len(sys.argv) > 2 else (sys.argv[1] if len(sys.argv) == 2 else "demo_output.pptx")
    create_ppt(template=tpl, output=out)
