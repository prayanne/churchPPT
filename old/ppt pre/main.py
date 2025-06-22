# ppt_generator_demo.py
"""
Basic demo script that creates a new PowerPoint file based on an existing template.

Dependencies
------------
python-pptx :
    pip install python-pptx
Pillow (optional, only if you embed images) :
    pip install pillow

Usage (CLI)
-----------
python ppt_generator_demo.py \
    --template template.pptx \
    --output demo_output.pptx

The script will generate three demo slides to showcase how to work with the
provided helper functions. Adapt `build_demo_slides()` to feed your own
content or integrate this module into a larger pipeline.
"""
import argparse
from pathlib import Path
from typing import List, Dict, Any

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN


def add_title_and_body_slide(prs: Presentation, title: str, body: str) -> None:
    """Add a *Title and Content* slide to *prs* with *title* and *body* text."""
    # Assume layout 1 is *Title and Content* for most PPTX templates
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)

    title_placeholder = slide.shapes.title
    body_placeholder = slide.shapes.placeholders[1]

    title_placeholder.text = title

    tf = body_placeholder.text_frame
    tf.clear()  # remove any existing paragraph

    p = tf.paragraphs[0]
    p.text = body
    p.font.size = Pt(18)
    p.alignment = PP_ALIGN.LEFT


def add_title_body_image_slide(
    prs: Presentation,
    title: str,
    body: str,
    image_path: Path,
    image_width_inch: float = 3.5,
) -> None:
    """Add a slide with title, body on the left and an image on the right."""
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)

    # --- Title & body on left ---
    title_placeholder = slide.shapes.title
    body_placeholder = slide.shapes.placeholders[1]

    title_placeholder.text = title

    tf = body_placeholder.text_frame
    tf.clear()
    tf.paragraphs[0].text = body
    tf.paragraphs[0].font.size = Pt(18)

    # --- Image on right side ---
    if image_path.exists():
        left = Inches(6)  # adjust as needed based on template dimensions
        top = Inches(2)
        slide.shapes.add_picture(str(image_path), left, top, width=Inches(image_width_inch))
    else:
        print(f"[WARN] Image not found: {image_path}. Skipping image insertion.")


def build_demo_slides(prs: Presentation) -> None:
    """Create a few demo slides showcasing helper utilities."""
    add_title_and_body_slide(
        prs,
        title="Demo Slide 1 – Title & Body",
        body=(
            "This slide demonstrates how to add a standard 'Title and Content' layout "
            "programmatically.\n\nFeel free to modify the helper function to change font "
            "sizes, colors, or alignment."
        ),
    )

    add_title_body_image_slide(
        prs,
        title="Demo Slide 2 – With Image",
        body="Here we add an image to the right while keeping text on the left side.",
        image_path=Path("sample_image.jpg"),
    )

    add_title_and_body_slide(
        prs,
        title="Next Steps",
        body="1. Replace demo slides with your real content.\n2. Define new layouts.\n3. Iterate and refine.",
    )


def create_ppt_from_template(template_path: Path, output_path: Path) -> None:
    """Load *template_path* and write a new presentation to *output_path*."""
    if not template_path.exists():
        raise FileNotFoundError(f"Template not found: {template_path}")

    prs = Presentation(str(template_path))

    # Option 1: keep template's cover slide, then append new slides
    # Option 2: clear all default slides created by template
    # Uncomment line below if you want to start fresh:
    # while prs.slides: prs.slides.remove(prs.slides[0])

    build_demo_slides(prs)

    prs.save(str(output_path))
    print(f"[INFO] Presentation saved to {output_path}")


# -------------------- CLI ENTRY POINT --------------------
if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="Generate a PPTX from a template.")
    parser.add_argument(
        "--template",
        type=Path,
        required=True,
        help="Path to the .pptx template file.",
    )
    parser.add_argument(
        "--output",
        type=Path,
        default=Path("generated_presentation.pptx"),
        help="Destination path for the generated presentation.",
    )
    args = parser.parse_args()

    create_ppt_from_template(args.template, args.output)
